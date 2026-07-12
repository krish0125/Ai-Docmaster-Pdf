"""PDF Conversion Service — Phase 2.

Converts between PDF and: Word, Excel, PowerPoint, Images, HTML, Text, EPUB.
All converters are wrapped in try/except with graceful ImportError handling so
the server still starts even if a library or system binary is missing.
"""

from __future__ import annotations
import os
import uuid
import tempfile


# ---------------------------------------------------------------------------
# Helper
# ---------------------------------------------------------------------------

def _save_output(content: bytes, suffix: str, upload_folder: str) -> tuple[str, str]:
    """Write *content* to a uuid-named file in *upload_folder*.
    Returns ``(filename, file_path)``.
    """
    fname = f"{uuid.uuid4().hex}{suffix}"
    path  = os.path.join(upload_folder, fname)
    os.makedirs(upload_folder, exist_ok=True)
    with open(path, 'wb') as f:
        f.write(content)
    return fname, path


# ───────────────────────────────────────────────────────────────────────────
# PDF ↔ Word (DOCX)
# ───────────────────────────────────────────────────────────────────────────

def pdf_to_word(pdf_path: str, upload_folder: str) -> tuple[str, str]:
    """Convert a PDF to a DOCX using pdf2docx.
    Returns ``(filename, file_path)``.
    """
    try:
        from pdf2docx import Converter
    except ImportError:
        raise RuntimeError("pdf2docx is not installed. Run: pip install pdf2docx")

    output_name = f"converted_{uuid.uuid4().hex}.docx"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)

    cv = Converter(pdf_path)
    cv.convert(output_path, start=0, end=None)
    cv.close()
    return output_name, output_path


def word_to_pdf(docx_path: str, upload_folder: str) -> tuple[str, str]:
    """Convert a DOCX to PDF.

    Tries docx2pdf first (requires Word on Windows), then falls back to
    LibreOffice headless. Raises RuntimeError if neither is available.
    """
    output_name = f"converted_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)

    # Attempt 1: docx2pdf (Word must be installed)
    try:
        from docx2pdf import convert
        convert(docx_path, output_path)
        if os.path.isfile(output_path):
            return output_name, output_path
    except Exception as e:
        print(f"[ConvertService] docx2pdf failed ({e}), trying LibreOffice...")

    # Attempt 2: LibreOffice headless
    import subprocess
    try:
        result = subprocess.run(
            ['soffice', '--headless', '--convert-to', 'pdf', '--outdir',
             upload_folder, docx_path],
            capture_output=True, timeout=60
        )
        if result.returncode != 0:
            raise RuntimeError(
                "Word-to-PDF conversion failed. Install Microsoft Word (docx2pdf) "
                "or LibreOffice (soffice in PATH). "
                f"LibreOffice error: {result.stderr.decode()[:200]}"
            )
    except Exception as e:
        raise RuntimeError(
            "Word-to-PDF conversion failed. Install Microsoft Word (docx2pdf) "
            "or LibreOffice (soffice in PATH). "
            f"LibreOffice details: {str(e)}"
        )

    # LibreOffice writes to the outdir with the same base name
    base = os.path.splitext(os.path.basename(docx_path))[0] + '.pdf'
    soffice_out = os.path.join(upload_folder, base)
    if os.path.isfile(soffice_out):
        os.replace(soffice_out, output_path)

    return output_name, output_path


# ───────────────────────────────────────────────────────────────────────────
# PDF ↔ Excel (XLSX)
# ───────────────────────────────────────────────────────────────────────────

def pdf_to_excel(pdf_path: str, upload_folder: str) -> tuple[str, str]:
    """Extract tables from a PDF into an XLSX workbook using pdfplumber + openpyxl."""
    try:
        import pdfplumber
        import openpyxl
    except ImportError as e:
        raise RuntimeError(f"Missing dependency: {e}. Run: pip install pdfplumber openpyxl")

    wb = openpyxl.Workbook()
    wb.remove(wb.active)          # remove default empty sheet

    with pdfplumber.open(pdf_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            tables = page.extract_tables()
            if not tables:
                # Fallback: dump raw text rows
                text = page.extract_text() or ''
                ws = wb.create_sheet(title=f"Page {page_num}")
                for row in text.split('\n'):
                    ws.append([row])
                continue
            for tbl_idx, table in enumerate(tables, start=1):
                title = f"P{page_num}_T{tbl_idx}"
                ws = wb.create_sheet(title=title[:31])   # max 31 chars
                for row in table:
                    ws.append([str(c) if c is not None else '' for c in row])

    if not wb.sheetnames:
        ws = wb.create_sheet("Data")
        ws.append(["No tables found in this PDF"])

    output_name = f"converted_{uuid.uuid4().hex}.xlsx"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)
    wb.save(output_path)
    return output_name, output_path


def excel_to_pdf(xlsx_path: str, upload_folder: str) -> tuple[str, str]:
    """Convert an XLSX to PDF using openpyxl (read) + reportlab (write)."""
    try:
        import openpyxl
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
        from reportlab.lib import colors
        from reportlab.lib.styles import getSampleStyleSheet
    except ImportError as e:
        raise RuntimeError(f"Missing dependency: {e}. Run: pip install openpyxl reportlab")

    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    output_name = f"converted_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)

    story = []
    styles = getSampleStyleSheet()

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        story.append(Paragraph(f"Sheet: {sheet_name}", styles['Heading2']))
        story.append(Spacer(1, 8))

        data = []
        for row in ws.iter_rows(values_only=True):
            data.append([str(c) if c is not None else '' for c in row])

        if data:
            t = Table(data, repeatRows=1)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#6C63FF')),
                ('TEXTCOLOR',  (0, 0), (-1, 0), colors.white),
                ('FONTSIZE',   (0, 0), (-1, -1), 7),
                ('GRID',       (0, 0), (-1, -1), 0.5, colors.grey),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#F5F3FF')]),
            ]))
            story.append(t)
        story.append(Spacer(1, 16))

    doc = SimpleDocTemplate(output_path, pagesize=A4)
    doc.build(story)
    return output_name, output_path


# ───────────────────────────────────────────────────────────────────────────
# PDF ↔ PowerPoint (PPTX)
# ───────────────────────────────────────────────────────────────────────────

def pdf_to_pptx(pdf_path: str, upload_folder: str) -> tuple[str, str]:
    """Convert PDF pages to a PPTX presentation.

    Extracts text from each page (pdfplumber) and creates one slide per page.
    Images from the PDF are not embedded (that requires pdf2image + poppler;
    if available they will be attempted).
    """
    try:
        import pdfplumber
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError as e:
        raise RuntimeError(f"Missing dependency: {e}. Run: pip install python-pptx pdfplumber")

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]   # blank

    with pdfplumber.open(pdf_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            slide = prs.slides.add_slide(blank_layout)
            text = page.extract_text() or f"[Page {page_num} — no extractable text]"

            # Title box at top
            txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = f"Page {page_num}"
            tf.paragraphs[0].runs[0].font.bold = True
            tf.paragraphs[0].runs[0].font.size = Pt(14)
            tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x6C, 0x63, 0xFF)

            # Content box
            cBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.8), Inches(9), Inches(6))
            cf = cBox.text_frame
            cf.word_wrap = True
            cf.text = text[:3000]          # cap per slide

    output_name = f"converted_{uuid.uuid4().hex}.pptx"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)
    prs.save(output_path)
    return output_name, output_path


def pptx_to_pdf(pptx_path: str, upload_folder: str) -> tuple[str, str]:
    """Convert a PPTX to PDF by extracting text/notes and rendering via reportlab."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib import colors
    except ImportError as e:
        raise RuntimeError(f"Missing dependency: {e}. Run: pip install python-pptx reportlab")

    prs = Presentation(pptx_path)
    styles = getSampleStyleSheet()
    slide_style = ParagraphStyle('Slide', parent=styles['Heading2'],
                                 textColor=colors.HexColor('#6C63FF'))
    body_style  = styles['BodyText']

    story = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        story.append(Paragraph(f"Slide {slide_num}", slide_style))
        story.append(Spacer(1, 4))
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        story.append(Paragraph(text, body_style))
        story.append(HRFlowable(width="100%", thickness=0.5,
                                 color=colors.HexColor('#E5E7EB')))
        story.append(Spacer(1, 12))

    output_name = f"converted_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)
    doc = SimpleDocTemplate(output_path, pagesize=landscape(A4))
    doc.build(story)
    return output_name, output_path


# ───────────────────────────────────────────────────────────────────────────
# PDF ↔ Images (JPG/PNG)
# ───────────────────────────────────────────────────────────────────────────

def pdf_to_images(pdf_path: str, upload_folder: str,
                  fmt: str = 'jpg', dpi: int = 150) -> list[tuple[str, str]]:
    """Convert each PDF page to an image file.

    Tries pdf2image (requires poppler) first, then falls back to PyMuPDF (fitz).
    Returns list of ``(filename, file_path)`` tuples, one per page.
    Set POPPLER_PATH in .env if poppler is not on PATH.
    """
    from config import Config
    poppler_path = getattr(Config, 'POPPLER_PATH', None) or os.environ.get('POPPLER_PATH')

    # Try pdf2image first (higher quality, needs poppler)
    try:
        from pdf2image import convert_from_path

        kwargs: dict = {'dpi': dpi}
        if poppler_path:
            kwargs['poppler_path'] = poppler_path

        images = convert_from_path(pdf_path, **kwargs)
        os.makedirs(upload_folder, exist_ok=True)

        results = []
        ext = 'jpg' if fmt.lower() in ('jpg', 'jpeg') else 'png'
        pil_fmt = 'JPEG' if ext == 'jpg' else 'PNG'

        for i, img in enumerate(images, start=1):
            fname = f"page_{i}_{uuid.uuid4().hex}.{ext}"
            path  = os.path.join(upload_folder, fname)
            img.save(path, pil_fmt)
            results.append((fname, path))

        return results

    except Exception:
        pass  # always fall through to fitz if pdf2image fails for any reason


    # Fallback: PyMuPDF (fitz)
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise RuntimeError(
            "pdf2image requires Poppler (not found on PATH). "
            "Install Poppler from https://github.com/oschwartz10612/poppler-windows/releases "
            "or set POPPLER_PATH in .env"
        )

    os.makedirs(upload_folder, exist_ok=True)
    doc = fitz.open(pdf_path)
    results = []
    zoom = dpi / 72.0  # 72 dpi is the base
    mat = fitz.Matrix(zoom, zoom)
    ext = 'jpg' if fmt.lower() in ('jpg', 'jpeg') else 'png'

    for i, page in enumerate(doc, start=1):
        pix = page.get_pixmap(matrix=mat)
        fname = f"page_{i}_{uuid.uuid4().hex}.{ext}"
        path  = os.path.join(upload_folder, fname)
        if ext == 'jpg':
            pix.save(path, 'jpg')
        else:
            pix.save(path, 'png')
        results.append((fname, path))

    doc.close()
    return results


def images_to_pdf(image_paths: list[str], upload_folder: str) -> tuple[str, str]:
    """Combine a list of images into a single PDF using img2pdf (with Pillow fallback)."""
    output_name = f"converted_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)

    # Try img2pdf first (preferred)
    try:
        import img2pdf
        with open(output_path, 'wb') as f:
            f.write(img2pdf.convert(image_paths))
        return output_name, output_path
    except Exception:
        pass  # fall through to Pillow

    # Fallback: Pillow
    from PIL import Image
    images_pil = []
    for p in image_paths:
        img = Image.open(p).convert('RGB')
        images_pil.append(img)

    if not images_pil:
        raise RuntimeError("No valid images found to convert.")

    first = images_pil[0]
    rest  = images_pil[1:] if len(images_pil) > 1 else []
    first.save(output_path, 'PDF', save_all=True, append_images=rest)

    return output_name, output_path


# ───────────────────────────────────────────────────────────────────────────
# PDF ↔ HTML
# ───────────────────────────────────────────────────────────────────────────

def pdf_to_html(pdf_path: str, upload_folder: str) -> tuple[str, str]:
    """Convert PDF to HTML by extracting text with pdfplumber and
    wrapping each page in a <section>."""
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("pdfplumber is not installed. Run: pip install pdfplumber")

    pages_html = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ''
            escaped = (text
                       .replace('&', '&amp;')
                       .replace('<', '&lt;')
                       .replace('>', '&gt;'))
            paragraphs = ''.join(f'<p>{line}</p>' for line in escaped.split('\n') if line.strip())
            pages_html.append(f'<section class="page"><h2>Page {i}</h2>{paragraphs}</section>')

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>Converted PDF</title>
<style>
  body{{font-family:Georgia,serif;max-width:800px;margin:2rem auto;color:#111;line-height:1.7}}
  .page{{border:1px solid #ddd;border-radius:8px;padding:1.5rem;margin:1.5rem 0}}
  h2{{font-size:1rem;color:#6C63FF;margin:0 0 .75rem}}
  p{{margin:.4rem 0}}
</style>
</head>
<body>
{''.join(pages_html)}
</body>
</html>"""

    output_name = f"converted_{uuid.uuid4().hex}.html"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html)
    return output_name, output_path


def html_to_pdf(html_path: str, upload_folder: str) -> tuple[str, str]:
    """Convert HTML to PDF.

    Tries weasyprint first; falls back to pdfkit+wkhtmltopdf.
    Raises RuntimeError with install instructions if neither is available.
    """
    output_name = f"converted_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)

    # Attempt 1: weasyprint
    try:
        from weasyprint import HTML
        HTML(filename=html_path).write_pdf(output_path)
        if os.path.isfile(output_path):
            return output_name, output_path
    except ImportError:
        pass
    except Exception as e:
        print(f"[ConvertService] weasyprint failed: {e}")

    # Attempt 2: pdfkit
    try:
        import pdfkit
        pdfkit.from_file(html_path, output_path)
        if os.path.isfile(output_path):
            return output_name, output_path
    except ImportError:
        pass
    except Exception as e:
        print(f"[ConvertService] pdfkit failed: {e}")

    raise RuntimeError(
        "HTML-to-PDF conversion failed. Install weasyprint (`pip install weasyprint`) "
        "or pdfkit+wkhtmltopdf (`pip install pdfkit` + https://wkhtmltopdf.org)."
    )


# ───────────────────────────────────────────────────────────────────────────
# PDF ↔ Text
# ───────────────────────────────────────────────────────────────────────────

def pdf_to_text(pdf_path: str, upload_folder: str) -> tuple[str, str]:
    """Extract all text from a PDF and save as .txt."""
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("pdfplumber is not installed.")

    lines = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            lines.append(f"=== PAGE {i} ===")
            text = page.extract_text() or '[No text on this page]'
            lines.append(text)
            lines.append('')

    output_name = f"converted_{uuid.uuid4().hex}.txt"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))
    return output_name, output_path


def text_to_pdf(txt_path: str, upload_folder: str) -> tuple[str, str]:
    """Convert a plain text file to PDF using reportlab."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
    except ImportError:
        raise RuntimeError("reportlab is not installed. Run: pip install reportlab")

    with open(txt_path, 'r', encoding='utf-8', errors='replace') as f:
        content = f.read()

    styles = getSampleStyleSheet()
    story = []
    for line in content.split('\n'):
        if line.strip():
            story.append(Paragraph(line.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;'),
                                   styles['BodyText']))
        else:
            story.append(Spacer(1, 6))

    output_name = f"converted_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)
    doc = SimpleDocTemplate(output_path, pagesize=A4)
    doc.build(story)
    return output_name, output_path


# ───────────────────────────────────────────────────────────────────────────
# PDF ↔ EPUB
# ───────────────────────────────────────────────────────────────────────────

def pdf_to_epub(pdf_path: str, upload_folder: str) -> tuple[str, str]:
    """Convert PDF to EPUB using pdfplumber (text) + ebooklib."""
    try:
        import pdfplumber
        import ebooklib
        from ebooklib import epub
    except ImportError as e:
        raise RuntimeError(f"Missing dependency: {e}. Run: pip install pdfplumber ebooklib")

    book = epub.EpubBook()
    base_name = os.path.splitext(os.path.basename(pdf_path))[0]
    book.set_identifier(uuid.uuid4().hex)
    book.set_title(base_name)
    book.set_language('en')

    chapters = []
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ''
            escaped = (text
                       .replace('&', '&amp;')
                       .replace('<', '&lt;')
                       .replace('>', '&gt;'))
            paragraphs = ''.join(
                f'<p>{line}</p>' for line in escaped.split('\n') if line.strip()
            )
            c = epub.EpubHtml(title=f'Page {i}', file_name=f'page_{i}.xhtml', lang='en')
            c.content = f'<html><body><h2>Page {i}</h2>{paragraphs}</body></html>'
            book.add_item(c)
            chapters.append(c)

    book.toc = tuple(chapters)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ['nav'] + chapters

    output_name = f"converted_{uuid.uuid4().hex}.epub"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)
    epub.write_epub(output_path, book)
    return output_name, output_path


def epub_to_pdf(epub_path: str, upload_folder: str) -> tuple[str, str]:
    """Convert EPUB to PDF by extracting text via ebooklib + rendering via reportlab."""
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
    except ImportError as e:
        raise RuntimeError(
            f"Missing dependency: {e}. "
            "Run: pip install ebooklib reportlab beautifulsoup4"
        )

    book = epub.read_epub(epub_path)
    styles = getSampleStyleSheet()
    story = []

    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        soup = BeautifulSoup(item.get_content(), 'html.parser')
        for tag in soup.find_all(['h1', 'h2', 'h3', 'p']):
            text = tag.get_text(strip=True)
            if not text:
                continue
            style = styles['Heading2'] if tag.name in ('h1', 'h2', 'h3') else styles['BodyText']
            story.append(Paragraph(
                text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;'),
                style
            ))
            story.append(Spacer(1, 4))

    output_name = f"converted_{uuid.uuid4().hex}.pdf"
    output_path = os.path.join(upload_folder, output_name)
    os.makedirs(upload_folder, exist_ok=True)
    doc = SimpleDocTemplate(output_path, pagesize=A4)
    doc.build(story if story else [Paragraph("No readable content found.", styles['BodyText'])])
    return output_name, output_path
