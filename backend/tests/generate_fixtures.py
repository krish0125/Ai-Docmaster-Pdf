# generate_fixtures.py
import os

fixtures_dir = os.path.join(os.path.dirname(__file__), 'fixtures')
os.makedirs(fixtures_dir, exist_ok=True)

# 1. Generate PDFs using reportlab
try:
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, PageBreak, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    print("Generating sample PDFs...")

    def make_pdf(filename, pages_count):
        path = os.path.join(fixtures_dir, filename)
        doc = SimpleDocTemplate(path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []
        for page in range(1, pages_count + 1):
            story.append(Paragraph(f"<h1>Sample PDF Page {page}</h1>", styles['Heading1']))
            story.append(Spacer(1, 20))
            story.append(Paragraph(
                f"This is some sample text on page {page} of {pages_count} pages in this document. "
                "We use this document to test AI DocMaster's PDF management, OCR text extraction, "
                "summarization, and metadata indexing tools.", styles['BodyText']
            ))
            if page < pages_count:
                story.append(PageBreak())
        doc.build(story)
        print(f"Generated: {filename} ({pages_count} pages)")

    make_pdf('sample_1page.pdf', 1)
    make_pdf('sample_5page.pdf', 5)
    make_pdf('sample_20page.pdf', 20)
except Exception as e:
    print(f"Error generating PDFs: {e}")

# 2. Generate DOCX
try:
    import docx
    print("Generating DOCX...")
    doc = docx.Document()
    doc.add_heading('Sample Word Document', 0)
    doc.add_paragraph('This is a simple sample DOCX file used to test the Word-to-PDF conversion service.')
    doc.save(os.path.join(fixtures_dir, 'sample.docx'))
    print("Generated: sample.docx")
except Exception as e:
    print(f"Error generating DOCX: {e}")

# 3. Generate XLSX
try:
    import openpyxl
    print("Generating XLSX...")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sales Data"
    ws.append(["Item", "Quantity", "Price", "Total"])
    ws.append(["Apples", 10, 1.50, 15.00])
    ws.append(["Oranges", 5, 2.00, 10.00])
    ws.append(["Bananas", 8, 0.50, 4.00])
    wb.save(os.path.join(fixtures_dir, 'sample.xlsx'))
    print("Generated: sample.xlsx")
except Exception as e:
    print(f"Error generating XLSX: {e}")

# 4. Generate PPTX
try:
    from pptx import Presentation
    print("Generating PPTX...")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Sample Presentation"
    slide.placeholders[1].text = "Created programmatically for test verification."
    prs.save(os.path.join(fixtures_dir, 'sample.pptx'))
    print("Generated: sample.pptx")
except Exception as e:
    print(f"Error generating PPTX: {e}")

# 5. Generate Images using Pillow
try:
    from PIL import Image, ImageDraw
    print("Generating sample PNG/JPG...")
    
    def make_img(filename, color):
        img = Image.new('RGB', (400, 300), color=color)
        d = ImageDraw.Draw(img)
        d.text((10, 10), f"Sample {filename.upper()}", fill=(255, 255, 255))
        img.save(os.path.join(fixtures_dir, filename))
        print(f"Generated: {filename}")

    make_img('sample.png', (108, 99, 255)) # custom purple
    make_img('sample.jpg', (255, 99, 132)) # custom pink
except Exception as e:
    print(f"Error generating Images: {e}")

# 6. Generate Text File
try:
    print("Generating TXT...")
    with open(os.path.join(fixtures_dir, 'sample.txt'), 'w', encoding='utf-8') as f:
        f.write("This is a simple plain text sample file used to verify text-to-pdf conversions and parsing.")
    print("Generated: sample.txt")
except Exception as e:
    print(f"Error generating TXT: {e}")

# 7. Generate EPUB
try:
    from ebooklib import epub
    print("Generating EPUB...")
    book = epub.EpubBook()
    book.set_identifier('id123456')
    book.set_title('Sample EPUB')
    book.set_language('en')
    book.add_author('AI DocMaster Test')
    
    c1 = epub.EpubHtml(title='Introduction', file_name='intro.xhtml', lang='en')
    c1.content = '<html><body><h1>Introduction</h1><p>This is a sample EPUB ebook for verification testing.</p></body></html>'
    book.add_item(c1)
    
    book.toc = (c1,)
    book.spine = ['nav', c1]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    
    epub.write_epub(os.path.join(fixtures_dir, 'sample.epub'), book)
    print("Generated: sample.epub")
except Exception as e:
    print(f"Error generating EPUB: {e}")

print("✅ All test fixtures generated successfully!")
